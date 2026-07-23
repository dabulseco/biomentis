"""Knowledge base for the tutor layer.

Owns a per-session Chroma vector index over uploaded files (PDF, PPTX, DOCX,
TXT) and web pages. The index is persisted to `data/tutor_kb/<session_id>/`
so a session can be resumed later. Retrieval is exposed via `retrieve(query)`,
which returns LangChain `Document`s with citation metadata.

The class deliberately hides LangChain / Chroma specifics from the rest of
the tutor package — `InstructionGenerator`, `TutorChat`, and the UI all just
call `KnowledgeBase.retrieve(...)` and read `Document.page_content` /
`Document.metadata`.
"""

from __future__ import annotations

import os
import re
import shutil
from dataclasses import dataclass, field
from typing import Iterable

# Imports are deferred so the tutor package can be imported on a system that
# hasn't installed the optional KB dependencies yet (e.g. a research-only
# deployment that never enables the tutor).

_CHUNK_SIZE_TOKENS = 500
_CHUNK_OVERLAP_TOKENS = 80


@dataclass
class KBStats:
    """Lightweight stats view used by the UI sidebar."""

    sources: int = 0
    chunks: int = 0
    bytes_indexed: int = 0
    last_updated: str | None = None
    source_names: list[str] = field(default_factory=list)


class KnowledgeBase:
    """A per-session vector index over uploaded documents and URLs.

    Args:
        session_id: Unique identifier for this learning session. Different
            sessions get different on-disk directories.
        path: Root directory for KB storage. Defaults to `./data/tutor_kb`.
        embedding_model: Name of the sentence-transformers model. The default
            is a small, fast, CPU-friendly model that downloads once and
            caches locally.
    """

    def __init__(
        self,
        session_id: str,
        path: str = "./data/tutor_kb",
        embedding_model: str = "sentence-transformers/all-MiniLM-L6-v2",
    ) -> None:
        self.session_id = session_id
        self.path = os.path.abspath(os.path.join(path, session_id))
        self.raw_dir = os.path.join(self.path, "raw")
        self.index_dir = os.path.join(self.path, "index")
        os.makedirs(self.raw_dir, exist_ok=True)
        os.makedirs(self.index_dir, exist_ok=True)

        self.embedding_model_name = embedding_model
        self._embeddings = None
        self._vectorstore = None
        self._source_names: set[str] = set()
        self._kb_signature: str = ""  # for instruction-card caching

    # ---- lazy initialization ----------------------------------------------

    def _ensure_embeddings(self):
        if self._embeddings is None:
            # Prefer the new langchain-huggingface package; fall back to the
            # deprecated langchain_community import for installs that haven't
            # migrated yet. Either way, the embedding model is local and
            # CPU-only.
            try:
                from langchain_huggingface import HuggingFaceEmbeddings
            except ImportError:
                from langchain_community.embeddings import HuggingFaceEmbeddings

            self._embeddings = HuggingFaceEmbeddings(
                model_name=self.embedding_model_name,
                model_kwargs={"device": "cpu"},
                encode_kwargs={"normalize_embeddings": True},
            )
        return self._embeddings

    def _ensure_vectorstore(self):
        if self._vectorstore is None:
            # Same dual-import trick for Chroma. langchain-chroma is the new
            # home; langchain_community still works on older installs.
            try:
                from langchain_chroma import Chroma
            except ImportError:
                from langchain_community.vectorstores import Chroma

            # Cold-start: hydrate `_source_names` from whatever already lives
            # in the index. If the index doesn't exist yet, the in-memory set
            # stays empty until the first add. We must NOT re-hydrate after
            # adds — the in-memory set is the source of truth post-add.
            self._vectorstore = Chroma(
                collection_name=f"tutor_kb_{self.session_id}",
                embedding_function=self._ensure_embeddings(),
                persist_directory=self.index_dir,
            )
            try:
                existing = self._vectorstore.get().get("metadatas", []) or []
                if existing and not self._source_names:
                    self._source_names = {
                        m.get("source", "unknown")
                        for m in existing
                        if m.get("source")
                    }
            except Exception:
                # Index may not exist yet; that's fine.
                pass
        return self._vectorstore

    # ---- public API --------------------------------------------------------

    def add_files(self, files: Iterable[str]) -> int:
        """Index a batch of file paths. Returns the number of chunks added.

        Accepts .pdf, .pptx, .docx, .txt. Unsupported extensions are skipped
        with a printed warning rather than raising — the UI is interactive
        and we don't want one bad file to abort an upload batch.

        A missing-dependency error (e.g. pypdf not installed) IS propagated,
        so the UI can tell the user to `pip install pypdf` instead of
        silently reporting 0 chunks.
        """
        from langchain_core.documents import Document

        new_docs: list[Document] = []
        for file_path in files:
            try:
                docs = self._load_file(file_path)
            except ImportError as e:
                # Missing optional dep (e.g. pypdf). Surface the message
                # so the user can install it.
                raise RuntimeError(
                    f"Missing dependency to read {os.path.basename(file_path)}: {e}. "
                    f"Run: pip install pypdf python-pptx python-docx"
                ) from e
            except RuntimeError:
                # Re-raise loader-originated RuntimeErrors (e.g. pypdf
                # import failed) — they carry actionable install instructions.
                raise
            except Exception as e:
                # Per-file error (corrupt PDF, etc.). Skip and log.
                print(f"KB: skipped {file_path}: {e!r}")
                continue
            new_docs.extend(docs)
            self._source_names.add(os.path.basename(file_path))

        if not new_docs:
            return 0

        chunks = self._chunk(new_docs)
        try:
            vs = self._ensure_vectorstore()
        except Exception as e:
            raise RuntimeError(
                f"Could not load embedding model "
                f"({self.embedding_model_name}): {e!r}. "
                "Try: pip install sentence-transformers"
            ) from e
        vs.add_documents(chunks)
        try:
            vs.persist()
        except Exception:
            # Chroma 0.5+ auto-persists; older versions need the explicit call.
            pass
        self._refresh_signature()
        return len(chunks)

    def add_urls(self, urls: Iterable[str]) -> int:
        """Fetch and index a batch of URLs. Returns chunks added."""
        from langchain_core.documents import Document

        new_docs: list[Document] = []
        for url in urls:
            try:
                docs = self._load_url(url)
            except Exception as e:
                print(f"KB: skipped {url}: {e!r}")
                continue
            new_docs.extend(docs)
            self._source_names.add(url)

        if not new_docs:
            return 0

        chunks = self._chunk(new_docs)
        try:
            vs = self._ensure_vectorstore()
        except Exception as e:
            raise RuntimeError(
                f"Could not load embedding model "
                f"({self.embedding_model_name}): {e!r}. "
                "Try: pip install sentence-transformers"
            ) from e
        vs.add_documents(chunks)
        try:
            vs.persist()
        except Exception:
            pass
        self._refresh_signature()
        return len(chunks)

    def retrieve(self, query: str, k: int = 4) -> list:
        """Return the top-k most relevant chunks for a query.

        Uses max-marginal-relevance so the same source isn't returned 4×
        when one page is on-topic. Returns a list of `langchain_core.documents.Document`
        with `.page_content` and `.metadata` (source, page, chunk_id).
        """
        if not self._source_names and self._vectorstore is None:
            return []
        vs = self._ensure_vectorstore()
        try:
            return vs.as_retriever(search_type="mmr", search_kwargs={"k": k}).invoke(query)
        except Exception:
            # Empty index, or MMR not configured yet — fall back to similarity.
            return vs.similarity_search(query, k=k)

    def stats(self) -> KBStats:
        """Read-only stats for the sidebar panel.

        Cheap path when the KB is empty: don't try to load the embedding
        model, which avoids requiring sentence-transformers to be installed
        on research-only deployments that never enable the tutor.
        """
        if not self._source_names and self._vectorstore is None:
            return KBStats()

        vs = self._ensure_vectorstore()
        try:
            data = vs.get()
            metadatas = data.get("metadatas", []) or []
            chunks = len(metadatas)
        except Exception:
            chunks = 0
        try:
            bytes_indexed = sum(
                os.path.getsize(os.path.join(root, f))
                for root, _, files in os.walk(self.index_dir)
                for f in files
            )
        except Exception:
            bytes_indexed = 0
        last_updated = None
        try:
            last_updated = datetime.fromtimestamp(  # noqa: F821 (import below)
                os.path.getmtime(self.index_dir)
            ).isoformat(timespec="seconds")
        except Exception:
            pass
        return KBStats(
            sources=len(self._source_names),
            chunks=chunks,
            bytes_indexed=bytes_indexed,
            last_updated=last_updated,
            source_names=sorted(self._source_names),
        )

    def clear(self) -> None:
        """Wipe the index and raw files for this session."""
        shutil.rmtree(self.path, ignore_errors=True)
        os.makedirs(self.raw_dir, exist_ok=True)
        os.makedirs(self.index_dir, exist_ok=True)
        self._vectorstore = None
        self._source_names = set()
        self._kb_signature = ""

    def kb_signature(self) -> str:
        """A short hash representing the current KB contents, for caching.

        The instruction generator hashes (event, kb_signature) so re-running
        the agent against the same KB doesn't re-bill the LLM for the same
        card. Returns "" if the KB is empty.
        """
        return self._kb_signature

    # ---- internals ---------------------------------------------------------

    def _load_file(self, file_path: str) -> list:
        """Load a single file into LangChain Documents with metadata."""
        from langchain_core.documents import Document

        ext = os.path.splitext(file_path)[1].lower()
        basename = os.path.basename(file_path)
        # PDFs are handled by `_read_pdf_pages` for per-page metadata, so
        # dispatch to them up front and skip the generic `_read_file_text`
        # check (which intentionally returns "" for PDFs).
        if ext == ".pdf":
            return self._read_pdf_pages(file_path, basename)
        text = self._read_file_text(file_path, ext)
        if not text.strip():
            return []
        return [
            Document(
                page_content=text,
                metadata={"source": basename, "page": None, "chunk_id": 0},
            )
        ]

    @staticmethod
    def _read_file_text(file_path: str, ext: str) -> str:
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        if ext == ".pdf":
            return ""  # handled separately for per-page metadata
        if ext == ".pptx":
            from pptx import Presentation

            prs = Presentation(file_path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            return "\n\n".join(parts)
        if ext == ".docx":
            from docx import Document as DocxDocument

            doc = DocxDocument(file_path)
            return "\n\n".join(p.text for p in doc.paragraphs if p.text)
        raise ValueError(f"Unsupported file type: {ext}")

    def _read_pdf_pages(self, file_path: str, basename: str) -> list:
        from langchain_core.documents import Document

        try:
            from pypdf import PdfReader
        except ImportError as e:
            raise RuntimeError("pypdf is required for PDF KB ingestion") from e

        reader = PdfReader(file_path)
        docs = []
        for i, page in enumerate(reader.pages, 1):
            try:
                text = page.extract_text() or ""
            except Exception:
                text = ""
            if not text.strip():
                continue
            docs.append(
                Document(
                    page_content=text,
                    metadata={"source": basename, "page": i, "chunk_id": 0},
                )
            )
        return docs

    def _load_url(self, url: str) -> list:
        from langchain_core.documents import Document

        import requests
        from bs4 import BeautifulSoup

        try:
            resp = requests.get(url, timeout=15, headers={"User-Agent": "Biomentis-Tutor/0.1"})
            resp.raise_for_status()
        except Exception as e:
            raise RuntimeError(f"fetch failed: {e!r}") from e
        soup = BeautifulSoup(resp.text, "lxml")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        text = re.sub(r"\s+", " ", soup.get_text(separator=" ")).strip()
        if not text:
            return []
        return [
            Document(
                page_content=text,
                metadata={"source": url, "page": None, "chunk_id": 0},
            )
        ]

    def _chunk(self, docs: list) -> list:
        """Token-based sliding window chunker with metadata-preserving overlap.

        Tries tiktoken first (token-accurate), falls back to char/4 if
        tiktoken isn't installed.
        """
        try:
            import tiktoken

            enc = tiktoken.get_encoding("cl100k_base")

            def tokenize(s: str) -> list[int]:
                return enc.encode(s)

            def detokenize(tokens: list[int]) -> str:
                return enc.decode(tokens)

        except Exception:

            def tokenize(s: str) -> list[str]:
                return s.split()

            def detokenize(tokens: list[str]) -> str:
                return " ".join(tokens)

        out = []
        for doc in docs:
            tokens = tokenize(doc.page_content)
            if not tokens:
                continue
            step = _CHUNK_SIZE_TOKENS - _CHUNK_OVERLAP_TOKENS
            chunk_idx = 0
            for start in range(0, max(1, len(tokens) - _CHUNK_OVERLAP_TOKENS), step):
                end = min(start + _CHUNK_SIZE_TOKENS, len(tokens))
                chunk_tokens = tokens[start:end]
                if not chunk_tokens:
                    continue
                text = detokenize(chunk_tokens)
                meta = dict(doc.metadata)
                meta["chunk_id"] = chunk_idx
                out.append(type(doc)(page_content=text, metadata=meta))
                chunk_idx += 1
                if end == len(tokens):
                    break
        return out

    def _refresh_signature(self) -> None:
        import hashlib

        h = hashlib.sha1()
        for name in sorted(self._source_names):
            h.update(name.encode("utf-8"))
            h.update(b"\x00")
        self._kb_signature = h.hexdigest()[:12]


# Late import for `stats()` (avoids paying for it on KB-only paths).
from datetime import datetime  # noqa: E402
